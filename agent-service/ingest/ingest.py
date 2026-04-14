"""Walk the corpus, convert docx/pptx/md/txt to text, chunk, embed via Voyage, upsert to Qdrant."""
import json
import os
import sys
import time
import urllib.error
import urllib.request
import uuid
from pathlib import Path

import tiktoken
from docx import Document
from pptx import Presentation
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, PointStruct, VectorParams


CORPUS_DIR = Path(os.environ.get("CORPUS_DIR", "/corpus"))
QDRANT_URL = os.environ.get("QDRANT_URL", "http://qdrant:6333")
COLLECTION_NAME = os.environ.get("COLLECTION_NAME", "cdpi_corpus")
# EMBED_PROVIDER: "tei" (local, default) or "voyage"
EMBED_PROVIDER = os.environ.get("EMBED_PROVIDER", "tei").lower()
EMBED_URL = os.environ.get("EMBED_URL", "http://embeddings:80/v1/embeddings")
# bge-small-en-v1.5 is 384-dim; voyage-3.5 is 1024-dim.
EMBEDDING_DIM = int(os.environ.get("EMBEDDING_DIM", "384" if EMBED_PROVIDER == "tei" else "1024"))
VOYAGE_API_KEY = os.environ.get("VOYAGE_API_KEY", "")
VOYAGE_MODEL = os.environ.get("VOYAGE_MODEL", "voyage-3.5")
CHUNK_TOKENS = 500
CHUNK_OVERLAP = 80
# TEI handles larger batches easily; Voyage free tier prefers small.
EMBED_BATCH = int(os.environ.get("EMBED_BATCH", "32" if EMBED_PROVIDER == "tei" else "8"))
EMBED_SLEEP = float(os.environ.get("EMBED_SLEEP", "0" if EMBED_PROVIDER == "tei" else "0.3"))

qdrant = QdrantClient(url=QDRANT_URL)
encoding = tiktoken.get_encoding("cl100k_base")


def extract_docx(path: Path) -> str:
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)


def extract_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


EXTRACTORS = {
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".md": extract_text,
    ".txt": extract_text,
}


def chunk_tokens(text: str, size: int, overlap: int) -> list[str]:
    tokens = encoding.encode(text)
    if not tokens:
        return []
    chunks = []
    i = 0
    while i < len(tokens):
        chunk = tokens[i : i + size]
        chunks.append(encoding.decode(chunk))
        if i + size >= len(tokens):
            break
        i += size - overlap
    return chunks


def voyage_embed(texts: list[str], input_type: str = "document") -> list[list[float]]:
    """Call Voyage /v1/embeddings with retry on 429."""
    body = json.dumps({
        "model": VOYAGE_MODEL,
        "input": texts,
        "input_type": input_type,
    }).encode()
    req = urllib.request.Request(
        "https://api.voyageai.com/v1/embeddings",
        data=body,
        headers={
            "Authorization": f"Bearer {VOYAGE_API_KEY}",
            "Content-Type": "application/json",
        },
        method="POST",
    )
    backoff = 2.0
    for attempt in range(6):
        try:
            with urllib.request.urlopen(req, timeout=60) as r:
                payload = json.loads(r.read())
                return [d["embedding"] for d in payload["data"]]
        except urllib.error.HTTPError as e:
            if e.code == 429 and attempt < 5:
                print(f"  rate-limited, backing off {backoff:.1f}s", flush=True)
                time.sleep(backoff)
                backoff *= 2
                continue
            print(f"  voyage error {e.code}: {e.read().decode(errors='ignore')[:300]}", flush=True)
            raise
        except Exception as e:
            if attempt < 5:
                print(f"  transient error {e}, retrying", flush=True)
                time.sleep(backoff)
                backoff *= 2
                continue
            raise
    raise RuntimeError("voyage_embed exhausted retries")


def tei_embed(texts: list[str], input_type: str = "document") -> list[list[float]]:
    """Call local TEI OpenAI-compatible /v1/embeddings endpoint. Retries while TEI warms up."""
    body = json.dumps({"input": texts}).encode()
    req = urllib.request.Request(
        EMBED_URL,
        data=body,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    for attempt in range(8):
        try:
            with urllib.request.urlopen(req, timeout=120) as r:
                payload = json.loads(r.read())
                return [d["embedding"] for d in payload["data"]]
        except urllib.error.HTTPError as e:
            if 500 <= e.code < 600 and attempt < 7:
                wait = 3 * (attempt + 1)
                print(f"  TEI 5xx ({e.code}), waiting {wait}s", flush=True)
                time.sleep(wait)
                continue
            raise
        except urllib.error.URLError as e:
            if attempt < 7:
                wait = 5 * (attempt + 1)
                print(f"  TEI not reachable ({e.reason}), waiting {wait}s (model may still be loading)", flush=True)
                time.sleep(wait)
                continue
            raise
    raise RuntimeError("tei_embed exhausted retries")


def embed(texts: list[str], input_type: str = "document") -> list[list[float]]:
    if EMBED_PROVIDER == "voyage":
        return voyage_embed(texts, input_type)
    return tei_embed(texts, input_type)


def ensure_collection() -> None:
    existing = [c.name for c in qdrant.get_collections().collections]
    if COLLECTION_NAME in existing:
        info = qdrant.get_collection(COLLECTION_NAME)
        existing_dim = info.config.params.vectors.size
        if existing_dim != EMBEDDING_DIM:
            print(
                f"Collection {COLLECTION_NAME} has dim {existing_dim} but ingest "
                f"needs {EMBEDDING_DIM}. Recreating.",
                flush=True,
            )
            qdrant.delete_collection(COLLECTION_NAME)
        else:
            return
    qdrant.create_collection(
        collection_name=COLLECTION_NAME,
        vectors_config=VectorParams(size=EMBEDDING_DIM, distance=Distance.COSINE),
    )
    print(f"Created Qdrant collection: {COLLECTION_NAME} ({EMBEDDING_DIM} dim)")


def ingest_file(path: Path) -> int:
    ext = path.suffix.lower()
    if ext not in EXTRACTORS:
        return 0
    try:
        text = EXTRACTORS[ext](path)
    except Exception as e:
        print(f"  skipped ({e})", flush=True)
        return 0
    if not text.strip():
        return 0

    rel_path = path.relative_to(CORPUS_DIR)
    parts = rel_path.parts
    category = parts[0] if len(parts) > 1 else ""
    subcategory = parts[1] if len(parts) > 2 else ""

    chunks = chunk_tokens(text, CHUNK_TOKENS, CHUNK_OVERLAP)
    if not chunks:
        return 0

    points: list[PointStruct] = []
    for batch_start in range(0, len(chunks), EMBED_BATCH):
        batch = chunks[batch_start : batch_start + EMBED_BATCH]
        vectors = embed(batch, input_type="document")
        for offset, (chunk, vec) in enumerate(zip(batch, vectors)):
            chunk_index = batch_start + offset
            # Deterministic ID so re-running ingest upserts instead of duplicating.
            point_id = str(uuid.uuid5(uuid.NAMESPACE_URL, f"{rel_path}#{chunk_index}"))
            points.append(
                PointStruct(
                    id=point_id,
                    vector=vec,
                    payload={
                        "source": str(rel_path),
                        "filename": path.name,
                        "category": category,
                        "subcategory": subcategory,
                        "chunk_index": chunk_index,
                        "text": chunk,
                    },
                )
            )
        time.sleep(EMBED_SLEEP)

    if points:
        qdrant.upsert(collection_name=COLLECTION_NAME, points=points, wait=True)
    return len(points)


def main() -> None:
    if not CORPUS_DIR.exists():
        print(f"Corpus directory {CORPUS_DIR} does not exist.", file=sys.stderr)
        sys.exit(1)
    if EMBED_PROVIDER == "voyage" and not VOYAGE_API_KEY:
        print("EMBED_PROVIDER=voyage but VOYAGE_API_KEY not set.", file=sys.stderr)
        sys.exit(1)
    print(f"Embed provider: {EMBED_PROVIDER} (dim={EMBEDDING_DIM})", flush=True)

    ensure_collection()

    files = sorted(
        p for p in CORPUS_DIR.rglob("*")
        if p.is_file() and p.suffix.lower() in EXTRACTORS
    )
    print(f"Found {len(files)} ingestible files in {CORPUS_DIR}", flush=True)

    total_chunks = 0
    for path in files:
        rel = path.relative_to(CORPUS_DIR)
        print(f"Ingesting {rel}", flush=True)
        count = ingest_file(path)
        print(f"  -> {count} chunks", flush=True)
        total_chunks += count

    info = qdrant.get_collection(COLLECTION_NAME)
    print(
        f"Done. Chunks added this run: {total_chunks}. "
        f"Collection total points: {info.points_count}",
        flush=True,
    )


if __name__ == "__main__":
    main()
